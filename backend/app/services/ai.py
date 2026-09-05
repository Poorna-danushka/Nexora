from pathlib import Path
from zipfile import BadZipFile
import json
import logging
import re

import httpx
from docx import Document
from docx.opc.exceptions import PackageNotFoundError
from pptx import Presentation
from pptx.exc import PackageNotFoundError as PptxPackageNotFoundError
from pypdf import PdfReader
from pypdf.errors import PdfReadError

from app.core.config import (
    AI_MAX_INPUT_CHARS,
    GEMINI_API_KEY,
    GEMINI_MODEL,
    GEMINI_TIMEOUT_SECONDS,
)
from app.schemas.ai import GeneratedQuizResponse
from app.services.ai_usage import set_provider_usage

GEMINI_GENERATE_CONTENT_URL = (
    "https://generativelanguage.googleapis.com/v1beta/"
    "models/{model}:generateContent"
)
MAX_SOURCE_TEXT_LENGTH = 120_000
MAX_PROVIDER_ERROR_RESPONSE_LENGTH = 1_000

logger = logging.getLogger(__name__)


class AIServiceError(Exception):
    """Base error for failures while calling the configured AI provider."""


class AIConfigurationError(AIServiceError):
    """Raised when the server is not configured to use the AI provider."""


class AIProviderError(AIServiceError):
    """Raised when the AI provider cannot return a usable response."""


class AIInputError(AIServiceError):
    """Raised when source material cannot be extracted or is empty."""


def _sanitized_provider_response(response_text: str) -> str:
    """Return a bounded provider error snippet without credential-like values."""
    sanitized = re.sub(r"(?:sk-|AIza)[A-Za-z0-9_-]+", "[REDACTED]", response_text)
    sanitized = re.sub(
        r'(?i)("?(?:api[_-]?key|authorization)"?\s*[:=]\s*["\']?)[^,\s"\']+',
        r"\1[REDACTED]",
        sanitized,
    )
    return sanitized[:MAX_PROVIDER_ERROR_RESPONSE_LENGTH]


def _provider_error_category(status_code: int) -> str:
    if status_code == 401:
        return "authentication"
    if status_code == 403:
        return "permission"
    if status_code == 429:
        return "rate_limit_or_quota"
    if status_code in (400, 404):
        return "invalid_request_or_model"
    if 500 <= status_code <= 599:
        return "provider_server_error"
    return "unexpected_http_error"


def _validate_input_size(*values: str) -> None:
    if sum(len(value) for value in values) > AI_MAX_INPUT_CHARS:
        raise AIInputError("AI input exceeds the configured size limit.")


def extract_material_text(path: Path, content_type: str) -> str:
    try:
        if content_type == "text/plain":
            text = path.read_text(encoding="utf-8")
        elif content_type == "application/pdf":
            text = "\n".join(page.extract_text() or "" for page in PdfReader(path).pages)
        elif content_type.endswith("wordprocessingml.document"):
            text = "\n".join(paragraph.text for paragraph in Document(path).paragraphs)
        elif content_type.endswith("presentationml.presentation"):
            text = "\n".join(
                shape.text
                for slide in Presentation(path).slides
                for shape in slide.shapes
                if hasattr(shape, "text")
            )
        else:
            raise AIInputError("Unsupported material type.")
    except (
        OSError,
        UnicodeError,
        ValueError,
        BadZipFile,
        PackageNotFoundError,
        PptxPackageNotFoundError,
        PdfReadError,
    ) as exc:
        raise AIInputError("Unable to read study material.") from exc

    text = text.strip()
    if not text:
        raise AIInputError("Study material contains no readable text.")
    return text[:MAX_SOURCE_TEXT_LENGTH]


def _request_completion(messages: list[dict[str, str]]) -> str:
    if not GEMINI_API_KEY:
        raise AIConfigurationError("GEMINI_API_KEY is not configured.")

    system_instruction = "\n\n".join(
        message["content"] for message in messages if message["role"] == "system"
    )
    contents = [
        {
            "role": "model" if message["role"] == "assistant" else "user",
            "parts": [{"text": message["content"]}],
        }
        for message in messages
        if message["role"] != "system"
    ]

    try:
        response = httpx.post(
            GEMINI_GENERATE_CONTENT_URL.format(model=GEMINI_MODEL),
            headers={
                "x-goog-api-key": GEMINI_API_KEY,
                "Content-Type": "application/json",
            },
            json={
                "systemInstruction": {"parts": [{"text": system_instruction}]},
                "contents": contents,
                "generationConfig": {"temperature": 0.2},
            },
            timeout=GEMINI_TIMEOUT_SECONDS,
        )
        response.raise_for_status()
        data = response.json()
    except httpx.TimeoutException as exc:
        logger.error("AI provider timeout: %s", exc.__class__.__name__)
        raise AIProviderError("The AI provider timed out.") from exc
    except httpx.HTTPStatusError as exc:
        logger.error(
            "AI provider HTTP error: status=%s category=%s response=%s",
            exc.response.status_code,
            _provider_error_category(exc.response.status_code),
            _sanitized_provider_response(exc.response.text),
        )
        raise AIProviderError("The AI provider returned an error.") from exc
    except httpx.RequestError as exc:
        logger.error("AI provider connection error: %s", exc.__class__.__name__)
        raise AIProviderError("The AI provider could not be reached.") from exc
    except ValueError as exc:
        logger.error("AI provider returned malformed JSON: %s", exc.__class__.__name__)
        raise AIProviderError("The AI provider returned invalid data.") from exc

    if not isinstance(data, dict):
        raise AIProviderError("The AI provider returned invalid data.")
    provider_usage = data.get("usageMetadata")
    if isinstance(provider_usage, dict):
        input_tokens = provider_usage.get("promptTokenCount")
        output_tokens = provider_usage.get("candidatesTokenCount")
        if (
            isinstance(input_tokens, int)
            and not isinstance(input_tokens, bool)
            and input_tokens >= 0
            and isinstance(output_tokens, int)
            and not isinstance(output_tokens, bool)
            and output_tokens >= 0
        ):
            set_provider_usage(
                {"input_tokens": input_tokens, "output_tokens": output_tokens}
            )
    candidates = data.get("candidates")
    if not isinstance(candidates, list) or not candidates:
        raise AIProviderError("The AI provider returned no answer.")
    content = (
        candidates[0].get("content") if isinstance(candidates[0], dict) else None
    )
    parts = content.get("parts") if isinstance(content, dict) else None
    answer = (
        "".join(part.get("text", "") for part in parts if isinstance(part, dict))
        if isinstance(parts, list)
        else ""
    )
    if not answer.strip():
        raise AIProviderError("The AI provider returned an empty answer.")
    return answer.strip()


def answer_conversation(messages: list[dict[str, str]]) -> str:
    if not messages:
        raise AIInputError("Conversation must contain a message.")
    _validate_input_size(*(message["content"] for message in messages))
    return _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "You are Nexora's study assistant. Answer clearly and helpfully. "
                    "When a question is academic, explain the reasoning and use examples "
                    "when useful. Do not claim to have access to information not provided."
                ),
            },
            *messages,
        ]
    )


def summarize_note(title: str, content: str) -> str:
    if not content.strip():
        raise ValueError("Note content must not be blank.")
    _validate_input_size(title, content)
    return _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "Summarize the student's note accurately and concisely. "
                    "Use a short paragraph followed by up to five key points. "
                    "Do not add facts that are not present in the note."
                ),
            },
            {
                "role": "user",
                "content": f"Note title: {title}\n\nNote content:\n{content}",
            },
        ]
    )


def answer_material_question(filename: str, source_text: str, question: str) -> str:
    if not source_text.strip():
        raise AIInputError("Study material contains no readable text.")
    if not question.strip():
        raise ValueError("Question must not be blank.")
    _validate_input_size(filename, source_text, question)
    return _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "Answer the student's question using only the provided study "
                    "material. If the answer is not present, say so clearly."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Study material: {filename}\n\n"
                    f"Material text:\n{source_text[:MAX_SOURCE_TEXT_LENGTH]}\n\n"
                    f"Question: {question}"
                ),
            },
        ]
    )


def generate_study_plan(
    subjects: str,
    goals: str,
    sessions: str,
    days: int,
    minutes_per_day: int,
    priorities: str | None,
) -> str:
    _validate_input_size(subjects, goals, sessions, priorities or "")
    return _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "Create a realistic student study plan using only the supplied "
                    "context. Return compact, mobile-friendly Markdown only: one "
                    "heading per day in the form 'Day N: focus', followed by at most "
                    "two short bullets containing the duration and learning action. "
                    "Do not include an introduction, conclusion, horizontal rules, "
                    "courses, or deadlines that were not supplied."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Plan length: {days} days\n"
                    f"Daily study time: {minutes_per_day} minutes\n"
                    f"Student priorities: {priorities or 'No additional priorities'}\n\n"
                    f"Subjects:\n{subjects or 'No subjects selected'}\n\n"
                    f"Existing goals:\n{goals or 'No active goals'}\n\n"
                    f"Existing sessions:\n{sessions or 'No upcoming sessions'}"
                ),
            },
        ]
    )


def generate_quiz(
    source_context: str,
    question_count: int,
    topic: str | None,
) -> GeneratedQuizResponse:
    _validate_input_size(source_context, topic or "")
    raw_response = _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "Generate a multiple-choice quiz using only the supplied source. "
                    "Return valid JSON only with this shape: "
                    '{"title":"string","questions":[{"question":"string",'
                    '"options":["string","string"],"correct_answer":"string",'
                    '"explanation":"string"}]}. '
                    "The correct_answer must exactly match one option. "
                    "Do not include markdown or additional keys."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Question count: {question_count}\n"
                    f"Topic: {topic or 'General coverage'}\n\n"
                    f"Source material:\n{source_context[:MAX_SOURCE_TEXT_LENGTH]}"
                ),
            },
        ]
    )
    try:
        parsed = json.loads(raw_response)
        result = GeneratedQuizResponse.model_validate(parsed)
    except (json.JSONDecodeError, TypeError, ValueError) as exc:
        raise AIProviderError("The AI provider returned an invalid quiz.") from exc
    if len(result.questions) != question_count:
        raise AIProviderError("The AI provider returned the wrong question count.")
    return result


def generate_practice_question(
    source_context: str,
    topic: str | None,
) -> GeneratedQuizResponse:
    return generate_quiz(source_context, 1, topic)


def explain_quiz_question(
    quiz_title: str,
    prompt: str,
    options: list[str],
    correct_option: int,
) -> str:
    if not prompt.strip() or not options or correct_option >= len(options):
        raise AIInputError("Quiz question is invalid.")
    _validate_input_size(quiz_title, prompt, *options)
    return _request_completion(
        [
            {
                "role": "system",
                "content": (
                    "Explain the quiz question clearly for a student. Identify why "
                    "the correct option is correct and briefly distinguish it from "
                    "the other options. Use only the supplied question data."
                ),
            },
            {
                "role": "user",
                "content": (
                    f"Quiz: {quiz_title}\n"
                    f"Question: {prompt}\n"
                    f"Options: {options}\n"
                    f"Correct option: {options[correct_option]}"
                ),
            },
        ]
    )
